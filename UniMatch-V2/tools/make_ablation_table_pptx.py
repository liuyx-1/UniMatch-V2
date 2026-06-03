"""Generate a single PowerPoint slide containing the Endoscapes-Seg50
ablation table as a native editable table (NOT an image).

Run:
    pip install python-pptx
    python tools/make_ablation_table_pptx.py
    # → outputs ablation_table.pptx in the current directory
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = 'ablation_table.pptx'

# ---------- table data ----------
HEADER = [
    'Variant', 'mIoU', 'mDice', 'mAP', 'PA', 'fwIoU',
    'mIoU_h', 'mIoU_b', 'mIoU_t', 'gap_HT↓',
    'cystic_plate', 'calot_triangle', 'cystic_artery',
    'cystic_duct', 'gallbladder', 'tool',
]
ROWS = [
    ['base',       '47.32', '56.00', '60.02', '94.93', '91.09',
                   '87.58', '17.25', '16.75', '70.83',
                   '5.30',  '16.75', '14.18', '32.26', '78.03', '88.86'],
    ['+text',      '49.16', '58.43', '60.65', '94.76', '91.09',
                   '87.57', '18.66', '25.22', '62.36',
                   '4.50',  '25.22', '17.47', '34.02', '77.73', '89.19'],
    ['+bias',      '46.31', '55.94', '57.13', '94.32', '90.33',
                   '86.18', '19.49', '12.09', '74.08',
                   '4.69',  '12.09', '20.21', '33.56', '76.66', '86.72'],
    ['+boundary',  '47.66', '56.82', '60.82', '95.28', '91.60',
                   '88.87', '17.72', '17.88', '71.00',
                   '5.53',  '17.88', '14.26', '33.36', '80.18', '90.42'],
    ['+bnd+bias',  '47.50', '57.23', '60.73', '94.35', '90.78',
                   '87.29', '21.14', '12.48', '74.81',
                   '4.26',  '12.48', '26.51', '32.64', '77.68', '88.78'],
]

# Bold-highlight cells (zero-indexed row, col).  Bigger-is-better unless noted.
# Skipping 'gap_HT' (lower-is-better) and the 'mIoU' approximate column.
BOLD_CELLS = {
    (1, 2),                                  # +text   mDice
    (1, 8), (1, 11), (1, 13),                # +text   mIoU_t, calot_triangle, cystic_duct
    (3, 3), (3, 4), (3, 5), (3, 6),          # +boundary mAP / PA / fwIoU / mIoU_h
    (3, 14), (3, 15),                        # +boundary gallbladder, tool
    (4, 7), (4, 12),                         # +bnd+bias mIoU_b, cystic_artery
}

# ---------- build slide ----------
prs = Presentation()
prs.slide_width  = Inches(13.33)   # 16:9
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# title
tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.5))
p = tx.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = 'Endoscapes-Seg50 · per-variant ablation (%)'
r.font.name = 'Times New Roman'; r.font.size = Pt(20); r.font.bold = True

# table dimensions
n_rows = len(ROWS) + 1
n_cols = len(HEADER)
tbl_left, tbl_top = Inches(0.4), Inches(1.0)
tbl_w, tbl_h    = Inches(12.5), Inches(3.6)
table = slide.shapes.add_table(n_rows, n_cols, tbl_left, tbl_top, tbl_w, tbl_h).table

# Column widths: first col wider for variant name
col_widths_in = [1.4] + [0.72] * 9 + [0.96] * 6
total_in = sum(col_widths_in)
scale = (tbl_w / Inches(1)) / total_in
for j, w in enumerate(col_widths_in):
    table.columns[j].width = Inches(w * scale)

# Header
HEADER_BG = RGBColor(0x2C, 0x3E, 0x50)
HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
ZEBRA_BG  = RGBColor(0xF6, 0xF8, 0xFA)
HIGHLIGHT = RGBColor(0xB8, 0x86, 0x0B)   # bold gold for best

for j, txt in enumerate(HEADER):
    cell = table.cell(0, j)
    cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BG
    tf = cell.text_frame; tf.clear()
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = txt
    run.font.name = 'Times New Roman'; run.font.size = Pt(10)
    run.font.bold = True; run.font.color.rgb = HEADER_FG

# Data rows
for i, row in enumerate(ROWS):
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        # zebra background on even data rows (0-indexed data row -> i)
        if i % 2 == 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = ZEBRA_BG
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf = cell.text_frame; tf.clear()
        para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        run = para.add_run(); run.text = val
        run.font.name = 'Times New Roman'; run.font.size = Pt(10)
        # bold variant column
        if j == 0:
            run.font.bold = True
            para.alignment = PP_ALIGN.LEFT
        # bold + gold for best cells
        if (i, j) in BOLD_CELLS:
            run.font.bold = True; run.font.color.rgb = HIGHLIGHT

# Footnote
foot = slide.shapes.add_textbox(Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.5))
fp = foot.text_frame.paragraphs[0]
fp.alignment = PP_ALIGN.LEFT
fr = fp.add_run()
fr.text = ('mIoU is recomputed from the 6 foreground IoUs and an estimated '
           '95.7% background IoU; all other numbers are read from test_metrics.csv. '
           'gap_HT is lower-is-better. Bold gold = best per column.')
fr.font.name = 'Times New Roman'; fr.font.size = Pt(9); fr.font.italic = True
fr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

prs.save(OUT)
print(f'saved {OUT}')
